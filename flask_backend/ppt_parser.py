from pptx import Presentation
from datetime import datetime
from thumbnail_generator import thumbnail_generator
from egnyte_file_operations import Egnyte_File_Operations
from file_download_logger import file_download_logger
from user_feedback_logger import user_feedback_logger
import utils
from logger import log
import objectpath

logger = log.getLogger()

class ppt_parser():

    def parse(path_to_file, file_relative_path):
        """Return data and meta-data of a ppt file in JSON format
                Args:
                    path_to_file (str): path of file
                Return:
                    data and meta-data of a given file in JSON format
            """
        file_data = None
        try:
            config = utils.config_parser()
            file_path_separator = config.get('egnyte', 'file_path_separator')

            ppt = Presentation(path_to_file)
            date_format = "%Y-%m-%d %H:%M:%S"
            file_data = ppt_parser.parse_metadata(ppt, path_to_file, file_relative_path)
            doc_id = utils.generate_doc_id(file_data['source_path'], file_path_separator)
            thumbnail_image_name = doc_id
            larger_thumbnail_list, smaller_thumbnail_list = thumbnail_generator.generate_thumbnail_image(path_to_file, thumbnail_image_name)
            if larger_thumbnail_list == [] or smaller_thumbnail_list == []:
                return None
            file_data['slides'] = ppt_parser.parse_content(ppt, larger_thumbnail_list, smaller_thumbnail_list, doc_id)
            file_data['title'] = ppt_parser.extract_document_level_title(file_data['slides'],file_data['file_name'])
        except Exception as e:
            logger.error("Failed to open file %s due to error %s" % (path_to_file, str(e)))
            return None

        return file_data

    def parse_metadata(ppt, path_to_file, file_relative_path):
        file_data = {}
        date_format = "%Y-%m-%dT%H:%M:%S"
        # Get the meta-data properties of a file
        meta_data = ppt.core_properties
        file_data['file_name'] = path_to_file.split("/")[-1]
        config = utils.config_parser()
        file_path_separator = config.get('egnyte', 'file_path_separator')
        corpus_directory_path = config.get('egnyte', 'corpus_path')
        egnyte_file_path = corpus_directory_path + file_relative_path[path_to_file].replace(file_path_separator, '/')
        file_data['url'] = Egnyte_File_Operations.get_file_navigation_link(egnyte_file_path=egnyte_file_path)
        file_data['source_path'] = file_relative_path[path_to_file]
        file_data['doc_type'] = path_to_file.split(".")[-1]
        file_data['created_by'] = meta_data.author
        file_data['revision'] = meta_data.revision
        file_data['modified_by'] = meta_data.last_modified_by
        doc_id = utils.generate_doc_id(file_data['source_path'], file_path_separator)
        file_data['num_downloads'] = file_download_logger.get_download_count_for_document(doc_id)
        user_feedback = user_feedback_logger.get_feedback_count_for_document(doc_id)
        file_data['ratings'] = user_feedback[0]
        file_data['num_likes'] = user_feedback[1]
        file_data['num_dislikes'] = user_feedback[2]
        checksum = path_to_file + 'checksum'
        file_data['checksum'] = file_relative_path[checksum]
        parser_version = path_to_file + 'parser_version'
        file_data['parser_version'] = file_relative_path[parser_version]
        if meta_data.created is not None:
            file_data['created_time'] = meta_data.created
        if meta_data.modified is not None:
             file_data['modified_time'] = meta_data.modified
        file_data['indexing_time'] = datetime.strftime(datetime.now(), date_format)
        return file_data

    def parse_content(ppt, larger_thumbnail_list, smaller_thumbnail_list, doc_id):
        """Return contents of a ppt file
            Args:
                ppt (obj): Object of Presentation
            Return:
                data contents of a given file
        """
        # import pprint
        # pp = pprint.PrettyPrinter(indent=2)

        slides_dict, slides_dict['slides'], shapes_dict, para_dict, line_dict, segments_dict,segments_dict_page1 = {}, {}, {}, {}, {}, {}, []
        all_slides_data = []
        es_config = utils.config_parser()
        parse_hidden_slides = es_config.get('generic', 'parse_hidden_slides')
        hidden_slide_count = 0
        for slide_idx, slide in enumerate(ppt.slides):
            # This loop traverses slides from a ppt fix
            # Check whether slide is hidden and parse_hidden_slides set to false then that hidden slide should not be parsed.

            if slide._element.get("show") == str(0) and parse_hidden_slides == 'false':
                hidden_slide_count += 1
                logger.info("Slide number %s is hidden." % str((ppt.slides.index(slide) + 1)))
                continue
            else:
                shapes_dict['page_number'] = ppt.slides.index(slide) + 1
                shapes_dict['page_id'] = doc_id + "_" + str(shapes_dict['page_number'])
                try:
                    shapes_dict['thumbnail_large'] = larger_thumbnail_list[ppt.slides.index(slide)-hidden_slide_count]
                except:
                    shapes_dict['thumbnail_large'] = ''
                try:
                    shapes_dict['thumbnail_small'] = smaller_thumbnail_list[ppt.slides.index(slide)-hidden_slide_count]
                except:
                    shapes_dict['thumbnail_small'] = ''
                try:
                    title = slide.shapes.title.text.encode('ascii', 'ignore').decode("utf-8").strip()
                    shapes_dict['title'] = title.replace('\u000b', ' ')
                except:
                    shapes_dict['title'] = 'Title Unknown'
                shapes_dict['shapes'] = []
                segments_dict['segments'] = []

                for shape in slide.shapes:
                    # This loop traverses all the shapes on a slide
                    para_dict['paras'] = []
                    line_dict['lines'] = []
                    segments_dict['segments'] = []

                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            # This loop traverses paragraphs from a shape
                            for run in para.runs:
                                # This loop traverses segments in a paragraph
                                if run.text.strip() != '':
                                    # Call to get segments method will return list of segments
                                    # If page number is 1 and title is "Title Unknown", add segment data into segments_dict_page1. segments_dict_page1 will be used to extract title for page 1.
                                    if shapes_dict['page_number'] == 1 and shapes_dict['title'] == 'Title Unknown':
                                        segments_dict['segments'], segments_dict_page1 = ppt_parser.get_segments(run, para, segments_dict['segments'],segments_dict_page1)
                                    else:
                                        segments_dict['segments'] = ppt_parser.get_segments(run, para, segments_dict['segments'])

                        if len(segments_dict['segments']):
                            line_dict['lines'].append(segments_dict.copy())
                        if len(line_dict['lines']):
                            para_dict['paras'].append(line_dict.copy())
                    if len(para_dict['paras']):
                        shapes_dict['shapes'].append(para_dict.copy())
                # If title for page 1 is "Title unknown", extract the title from font_size of text.
                if shapes_dict['page_number'] == 1 and shapes_dict['title'] == 'Title Unknown' and segments_dict_page1:
                    first_page_title =  ppt_parser.extract_title(segments_dict_page1)
                    if first_page_title:
                        shapes_dict['title'] = first_page_title
                    else:
                        shapes_dict['title'] = ''
                all_slides_data.append(shapes_dict.copy())
        return all_slides_data

    def get_segments(run, para, segments,segments_page1=None):
        """Return segments from a slide. Segment if segments_dict['segments']:s are nothing but text separated by it's formatting
            Args:
                run (obj): Object of paragraph.runs
                para (obj): Object of text_frame
            Return:
                data contents of a given file
        """
        segment_dict={}
        text = run.text.strip()
        # record font bold property of a segment
        bold = run.font.bold
        if bold == None: bold = False
        # record font size property of a segment
        if run.font.size:
            font_size = run.font.size.pt
        else:
            font_size = run.font.size
        # record font Italic property of a segment
        italic = run.font.italic
        if italic == None: italic = False

        # record font color property of a segment based on its type
        color_format = run.font.color
        color = None
        if color_format.type:
            if str(color_format.type).find('RGB') != -1:
                color = color_format.rgb
            elif str(color_format.type).find('SCHEME') != -1:
                color = color_format.theme_color
        else:
            color = None
        # record indentation level of a segment
        bullet_level = para.level
        segment_dict = {"text": text, "indentation": bullet_level, "font_size": font_size, "bold": bold, "italic": italic, "color": color}
        segments.append(segment_dict)
        if segments_page1 is not None:
            segments_page1.append(segment_dict)
            return segments, segments_page1
        return segments

    def extract_title(data):
        font_size = []
        for dist in data:
            fontsize= dist['font_size']
            font_size.append(fontsize)
        if any(elem is None for elem in font_size):
            title = data[0].get('text')

        else:
            m = max(font_size)
            max_index = [i for i, j in enumerate(font_size) if j == m]
            title_ = []
            for index in max_index:
                text = data[index].get('text')
                title_.append(text)
            title =' '.join(title_)
        return title


    def extract_document_level_title(data, file_name):
        tree = objectpath.Tree(data)
        get_title = [title for title in tree.execute("$..*[@.page_number is 1].title")]
        if get_title:
            if get_title[0] == 'Title Unknown' or get_title[0] == '':
                return file_name.rsplit('.', 1)[0]
            else:
                return get_title[0]
        else :
            return file_name.rsplit('.', 1)[0]

# def __main():
#     import pprint
#     path_to_file = "D:/KMP/knowledgemanagementportal/kmt/flask_backend/temp_storage/BilleoXpressBuy.pptx"
#     file_data = ppt_parser.parse(path_to_file)
#     pp = pprint.PrettyPrinter(indent=2)
#     pp.pprint(file_data)
#
#     import json
#     json_file = open('json_output_file.json', 'w')
#     json.dump(file_data, json_file)
#
#
# if __name__ == '__main__':
#     __main()

